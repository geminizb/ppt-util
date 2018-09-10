from pptx import Presentation
import copy
import six

# interal used by duplicate_slide
def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]

# delete the specified index of slide
def delete_slide(pres, index):
    xml_slides = pres.slides._sldIdLst
    xml_slides.remove(list(xml_slides)[index])

# move the slide to the target index
def move_slide(pres, slide, index):
    old_index = pres.slides.index(slide)
    xml_slides = pres.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(index, slides[old_index])

# copy the specified index of slide to the target new_index
def duplicate_slide(pres, index, new_index):
    blank_slide_layout = _get_blank_slide_layout(pres)
    copied_slide = pres.slides.add_slide(blank_slide_layout)

    source_slide = pres.slides[index]
    for shp in source_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(source_slide.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                                    value._target,
                                                    value.rId)
    
    move_slide(pres, copied_slide, new_index)
    
    return copied_slide
